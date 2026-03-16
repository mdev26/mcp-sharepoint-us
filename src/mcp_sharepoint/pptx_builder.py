"""
PowerPoint Document Builder for MCP SharePoint Server

This module provides a PowerPointBuilder class for creating formatted PowerPoint presentations
programmatically. It supports slides with various layouts, text boxes, bullet points, images,
tables, shapes, and comprehensive formatting options.

Example usage:
    builder = PowerPointBuilder()
    content = [
        {
            "type": "slide",
            "layout": "title",
            "title": "My Presentation",
            "subtitle": "Created with python-pptx"
        },
        {
            "type": "slide",
            "layout": "content",
            "title": "Key Points",
            "content": [
                {"type": "bullets", "items": ["Point 1", "Point 2", "Point 3"]}
            ]
        }
    ]
    pptx_bytes = builder.build(content)
"""

import base64
import logging
import requests
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


# Named color mappings (21 common colors, consistent with Word)
NAMED_COLORS = {
    "black": (0, 0, 0),
    "white": (255, 255, 255),
    "red": (255, 0, 0),
    "green": (0, 128, 0),
    "blue": (0, 0, 255),
    "yellow": (255, 255, 0),
    "cyan": (0, 255, 255),
    "magenta": (255, 0, 255),
    "orange": (255, 165, 0),
    "purple": (128, 0, 128),
    "pink": (255, 192, 203),
    "brown": (165, 42, 42),
    "gray": (128, 128, 128),
    "grey": (128, 128, 128),  # Alias for gray
    "light_gray": (211, 211, 211),
    "dark_gray": (169, 169, 169),
    "navy": (0, 0, 128),
    "teal": (0, 128, 128),
    "lime": (0, 255, 0),
    "maroon": (128, 0, 0),
    "olive": (128, 128, 0),
}

# Alignment mappings
ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

# Vertical alignment mappings
VERTICAL_ALIGNMENTS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}

# Shape type mappings
SHAPE_TYPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "arrow_left": MSO_SHAPE.LEFT_ARROW,
    "star": MSO_SHAPE.STAR_5_POINT,
    "hexagon": MSO_SHAPE.HEXAGON,
    "cloud": MSO_SHAPE.CLOUD,
}


class PowerPointBuilder:
    """
    A builder class for creating formatted PowerPoint presentations.

    Features:
    - Multiple slide layouts (title, content, blank, section_header, etc.)
    - Text boxes with positioning and formatting
    - Bullet points (ordered and unordered lists)
    - Images from URLs or base64 data
    - Tables with customizable styling
    - Shapes (rectangles, circles, arrows, etc.)
    - Color parsing (named colors, RGB tuples, RGB strings)
    - Font customization (size, color, bold, italic)
    - Alignment and positioning control
    """

    def __init__(self):
        """Initialize a new PowerPoint presentation."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)  # Standard 4:3 ratio
        self.prs.slide_height = Inches(7.5)
        self.stats = {
            "slides": 0,
            "text_boxes": 0,
            "images": 0,
            "tables": 0,
            "shapes": 0,
            "bullets": 0,
        }

    def _parse_color(self, color: Any) -> Optional[RGBColor]:
        """
        Parse a color from various formats into RGBColor.

        Supports:
        - Named colors: "red", "blue", "green", etc.
        - RGB tuples/lists: [255, 0, 0] or (255, 0, 0)
        - RGB strings: "255,0,0" or "255, 0, 0"

        Args:
            color: Color specification in various formats

        Returns:
            RGBColor object or None if parsing fails
        """
        if color is None:
            return None

        # Handle named colors
        if isinstance(color, str):
            color_lower = color.lower().replace(" ", "_")
            if color_lower in NAMED_COLORS:
                r, g, b = NAMED_COLORS[color_lower]
                return RGBColor(r, g, b)

            # Try parsing as comma-separated RGB string
            if "," in color:
                try:
                    parts = [int(x.strip()) for x in color.split(",")]
                    if len(parts) == 3:
                        r, g, b = parts
                        return RGBColor(r, g, b)
                except (ValueError, TypeError):
                    pass

        # Handle RGB tuples/lists
        elif isinstance(color, (list, tuple)) and len(color) == 3:
            try:
                r, g, b = [int(c) for c in color]
                if 0 <= r <= 255 and 0 <= g <= 255 and 0 <= b <= 255:
                    return RGBColor(r, g, b)
            except (ValueError, TypeError):
                pass

        logger.warning(f"Could not parse color: {color}")
        return None

    def _fetch_image(self, source: str) -> Optional[BytesIO]:
        """
        Fetch image from URL or decode from base64.

        Supports:
        - HTTP/HTTPS URLs
        - Base64 encoded data
        - Data URIs (data:image/png;base64,...)

        Args:
            source: Image source (URL or base64 string)

        Returns:
            BytesIO containing image data or None if fetch fails
        """
        try:
            # Handle HTTP/HTTPS URLs
            if source.startswith(('http://', 'https://')):
                response = requests.get(source, timeout=30)
                response.raise_for_status()
                return BytesIO(response.content)

            # Handle data URIs
            elif source.startswith('data:image'):
                # Extract base64 data after 'base64,'
                if 'base64,' in source:
                    base64_data = source.split('base64,', 1)[1]
                    image_data = base64.b64decode(base64_data)
                    return BytesIO(image_data)

            # Handle raw base64
            else:
                try:
                    image_data = base64.b64decode(source)
                    return BytesIO(image_data)
                except Exception:
                    pass

        except Exception as e:
            logger.error(f"Failed to fetch image from {source[:100]}: {e}")
            return None

        return None

    def _add_slide(self, item: Dict[str, Any]) -> None:
        """
        Add a new slide with the specified layout and content.

        Supported layouts:
        - title: Title slide with title and subtitle
        - content: Content slide with title and content area
        - section_header: Section header with title and subtitle
        - blank: Blank slide with no predefined placeholders
        - title_only: Slide with only a title placeholder

        Args:
            item: Dictionary with slide configuration
        """
        layout_name = item.get("layout", "blank").lower()

        # Map layout names to indices (based on default PowerPoint template)
        layout_map = {
            "title": 0,           # Title Slide
            "content": 1,         # Title and Content
            "section_header": 2,  # Section Header
            "title_only": 5,      # Title Only
            "blank": 6,           # Blank
        }

        layout_idx = layout_map.get(layout_name, 6)  # Default to blank

        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
            slide = self.prs.slides.add_slide(slide_layout)
            self.stats["slides"] += 1

            # Handle title and subtitle for title/section_header layouts
            if layout_name in ["title", "section_header"]:
                title_text = item.get("title", "")
                subtitle_text = item.get("subtitle", "")

                if slide.shapes.title:
                    slide.shapes.title.text = title_text
                    self._format_text_frame(slide.shapes.title.text_frame, item.get("title_format", {}))

                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle_text
                    self._format_text_frame(slide.placeholders[1].text_frame, item.get("subtitle_format", {}))

            # Handle title for content/title_only layouts
            elif layout_name in ["content", "title_only"]:
                title_text = item.get("title", "")
                if slide.shapes.title:
                    slide.shapes.title.text = title_text
                    self._format_text_frame(slide.shapes.title.text_frame, item.get("title_format", {}))

            # Process content items for content layout or blank slides
            content_items = item.get("content", [])
            if content_items:
                # Store current slide for content processing
                self.current_slide = slide
                for content_item in content_items:
                    self._process_content_item(content_item)

        except Exception as e:
            logger.error(f"Failed to add slide: {e}")
            # Add error slide
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            text_box.text = f"Error creating slide: {str(e)}"

    def _format_text_frame(self, text_frame: Any, formatting: Dict[str, Any]) -> None:
        """
        Apply formatting to a text frame.

        Args:
            text_frame: The text frame to format
            formatting: Dictionary with formatting options (font_size, color, bold, italic, etc.)
        """
        if not formatting:
            return

        for paragraph in text_frame.paragraphs:
            # Paragraph alignment
            align = formatting.get("alignment", "").lower()
            if align in ALIGNMENTS:
                paragraph.alignment = ALIGNMENTS[align]

            for run in paragraph.runs:
                font = run.font

                # Font size
                if "font_size" in formatting:
                    font.size = Pt(formatting["font_size"])

                # Font color
                if "color" in formatting:
                    color = self._parse_color(formatting["color"])
                    if color:
                        font.color.rgb = color

                # Bold/Italic
                if "bold" in formatting:
                    font.bold = formatting["bold"]
                if "italic" in formatting:
                    font.italic = formatting["italic"]

        # Vertical alignment
        v_align = formatting.get("vertical_alignment", "").lower()
        if v_align in VERTICAL_ALIGNMENTS:
            text_frame.vertical_anchor = VERTICAL_ALIGNMENTS[v_align]

    def _process_content_item(self, item: Dict[str, Any]) -> None:
        """
        Process a content item and add it to the current slide.

        Args:
            item: Dictionary describing the content item
        """
        content_type = item.get("type", "").lower()

        if content_type == "text_box":
            self._add_text_box(item)
        elif content_type == "bullets":
            self._add_bullets(item)
        elif content_type == "image":
            self._add_image(item)
        elif content_type == "table":
            self._add_table(item)
        elif content_type == "shape":
            self._add_shape(item)
        else:
            logger.warning(f"Unknown content type: {content_type}")

    def _add_text_box(self, item: Dict[str, Any]) -> None:
        """
        Add a text box to the current slide.

        Args:
            item: Dictionary with text box configuration (text, left, top, width, height, formatting)
        """
        try:
            text = item.get("text", "")
            left = Inches(item.get("left", 1))
            top = Inches(item.get("top", 2))
            width = Inches(item.get("width", 8))
            height = Inches(item.get("height", 1))

            text_box = self.current_slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = text

            # Apply formatting
            formatting = item.get("formatting", {})
            self._format_text_frame(text_frame, formatting)

            self.stats["text_boxes"] += 1

        except Exception as e:
            logger.error(f"Failed to add text box: {e}")

    def _add_bullets(self, item: Dict[str, Any]) -> None:
        """
        Add a bulleted or numbered list to the current slide.

        Args:
            item: Dictionary with bullets configuration (items, left, top, width, height, ordered, font_size)
        """
        try:
            items = item.get("items", [])
            left = Inches(item.get("left", 1))
            top = Inches(item.get("top", 2))
            width = Inches(item.get("width", 8))
            height = Inches(item.get("height", 4))
            ordered = item.get("ordered", False)
            font_size = item.get("font_size", 18)

            text_box = self.current_slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for idx, bullet_text in enumerate(items):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet_text
                p.level = 0

                # Enable bullet or numbering formatting
                if ordered:
                    # For numbered lists, use manual numbering for simplicity
                    # (PowerPoint's autonumbering via XML is complex)
                    p.text = f"{idx + 1}. {bullet_text}"

                # Set font size for the text via runs
                for run in p.runs:
                    run.font.size = Pt(font_size)

            # Apply additional formatting if specified
            formatting = item.get("formatting", {})
            if formatting:
                self._format_text_frame(text_frame, formatting)

            self.stats["bullets"] += len(items)

        except Exception as e:
            logger.error(f"Failed to add bullets: {e}")

    def _add_image(self, item: Dict[str, Any]) -> None:
        """
        Add an image to the current slide.

        Args:
            item: Dictionary with image configuration (source, left, top, width, height)
        """
        try:
            source = item.get("source", "")
            left = Inches(item.get("left", 1))
            top = Inches(item.get("top", 2))
            width = item.get("width")  # Optional
            height = item.get("height")  # Optional

            image_stream = self._fetch_image(source)
            if image_stream:
                if width and height:
                    self.current_slide.shapes.add_picture(
                        image_stream, left, top, width=Inches(width), height=Inches(height)
                    )
                elif width:
                    self.current_slide.shapes.add_picture(
                        image_stream, left, top, width=Inches(width)
                    )
                elif height:
                    self.current_slide.shapes.add_picture(
                        image_stream, left, top, height=Inches(height)
                    )
                else:
                    self.current_slide.shapes.add_picture(image_stream, left, top)

                self.stats["images"] += 1
            else:
                logger.error(f"Failed to fetch image from {source[:100]}")

        except Exception as e:
            logger.error(f"Failed to add image: {e}")

    def _add_table(self, item: Dict[str, Any]) -> None:
        """
        Add a table to the current slide.

        Args:
            item: Dictionary with table configuration (rows, left, top, width, height, header_color, cell_color)
        """
        try:
            rows_data = item.get("rows", [])
            if not rows_data:
                return

            rows = len(rows_data)
            cols = len(rows_data[0]) if rows_data else 0

            left = Inches(item.get("left", 1))
            top = Inches(item.get("top", 2))
            width = Inches(item.get("width", 8))
            height = Inches(item.get("height", 3))

            table = self.current_slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Fill table cells
            for row_idx, row_data in enumerate(rows_data):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.rows[row_idx].cells[col_idx]
                    cell.text = str(cell_data)

                    # Format header row
                    if row_idx == 0 and item.get("has_header", True):
                        header_color = self._parse_color(item.get("header_color", [68, 114, 196]))
                        if header_color:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = header_color

                        # Make header text white and bold
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True

            self.stats["tables"] += 1

        except Exception as e:
            logger.error(f"Failed to add table: {e}")

    def _add_shape(self, item: Dict[str, Any]) -> None:
        """
        Add a shape to the current slide.

        Args:
            item: Dictionary with shape configuration (shape_type, left, top, width, height, fill_color, line_color)
        """
        try:
            shape_type_name = item.get("shape_type", "rectangle").lower()
            shape_type = SHAPE_TYPES.get(shape_type_name, MSO_SHAPE.RECTANGLE)

            left = Inches(item.get("left", 1))
            top = Inches(item.get("top", 2))
            width = Inches(item.get("width", 2))
            height = Inches(item.get("height", 2))

            shape = self.current_slide.shapes.add_shape(shape_type, left, top, width, height)

            # Fill color
            fill_color = self._parse_color(item.get("fill_color"))
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
            elif item.get("fill_color") is None:
                # No fill specified, use default
                pass
            else:
                # Explicitly no fill
                shape.fill.background()

            # Line color
            line_color = self._parse_color(item.get("line_color"))
            if line_color:
                shape.line.color.rgb = line_color

            # Line width
            if "line_width" in item:
                shape.line.width = Pt(item["line_width"])

            # Add text to shape if specified
            if "text" in item:
                text_frame = shape.text_frame
                text_frame.text = item["text"]

                # Apply text formatting
                if "text_format" in item:
                    self._format_text_frame(text_frame, item["text_format"])

            self.stats["shapes"] += 1

        except Exception as e:
            logger.error(f"Failed to add shape: {e}")

    def build(self, content: List[Dict[str, Any]]) -> bytes:
        """
        Build a PowerPoint presentation from the content specification.

        Args:
            content: List of slide/content dictionaries

        Returns:
            Bytes of the generated .pptx file
        """
        try:
            # Process each item in the content array
            for item in content:
                item_type = item.get("type", "").lower()

                if item_type == "slide":
                    self._add_slide(item)
                else:
                    logger.warning(f"Unknown top-level type: {item_type}. Use 'slide' type.")

            # Save to BytesIO
            pptx_bytes = BytesIO()
            self.prs.save(pptx_bytes)
            pptx_bytes.seek(0)

            return pptx_bytes.getvalue()

        except Exception as e:
            logger.error(f"Failed to build PowerPoint presentation: {e}")
            raise

    def get_stats(self) -> Dict[str, int]:
        """
        Get statistics about the generated presentation.

        Returns:
            Dictionary with counts of various elements
        """
        return self.stats.copy()


def create_powerpoint(content: List[Dict[str, Any]]) -> bytes:
    """
    Convenience function to create a PowerPoint presentation.

    Args:
        content: List of slide/content dictionaries

    Returns:
        Bytes of the generated .pptx file
    """
    builder = PowerPointBuilder()
    return builder.build(content)
