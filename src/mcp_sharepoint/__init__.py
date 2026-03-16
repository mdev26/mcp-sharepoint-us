"""
SharePoint MCP Server with Modern Azure AD Authentication
"""
import os
import logging
import asyncio
from functools import wraps
from typing import Optional
import base64
import mimetypes
from io import BytesIO

from openpyxl import load_workbook

from mcp.server import Server
from mcp.types import Resource, Tool, TextContent, ImageContent, EmbeddedResource
from pydantic import AnyUrl
import mcp.server.stdio

from .graph_api import GraphAPIClient
from .docx_builder import create_word_document, edit_word_document
from .pptx_builder import create_powerpoint

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize MCP server
app = Server("mcp-sharepoint")

# Global Graph API client and authenticator
graph_client: Optional[GraphAPIClient] = None
authenticator = None


def ensure_context(func):
    """Decorator to ensure Graph API client is available"""
    @wraps(func)
    async def wrapper(*args, **kwargs):
        global graph_client, authenticator
        if graph_client is None:
            try:
                from .auth import SharePointAuthenticator

                # Get credentials
                site_url = os.getenv("SHP_SITE_URL")
                client_id = os.getenv("SHP_ID_APP")
                client_secret = os.getenv("SHP_ID_APP_SECRET")
                tenant_id = os.getenv("SHP_TENANT_ID")
                document_library = os.getenv("SHP_DOC_LIBRARY", "Shared Documents")
                cloud = "government" if ".sharepoint.us" in site_url else "commercial"

                # Create shared authenticator
                authenticator = SharePointAuthenticator(
                    site_url=site_url,
                    client_id=client_id,
                    client_secret=client_secret,
                    tenant_id=tenant_id,
                    cloud=cloud
                )

                # Create Graph API client with direct token access
                def get_token():
                    """Get access token for Graph API"""
                    return authenticator.get_access_token()

                graph_client = GraphAPIClient(
                    site_url=site_url,
                    token_callback=get_token,
                    document_library_name=document_library
                )
                logger.info(f"Connected to SharePoint: {site_url}")

            except Exception as e:
                logger.error(f"Failed to initialize Graph API client: {e}", exc_info=True)
                raise RuntimeError(
                    f"Graph API authentication failed: {e}. "
                    "Please check your environment variables and ensure:\n"
                    "1. SHP_TENANT_ID is set correctly\n"
                    "2. Your Azure AD app has Microsoft Graph API permissions\n"
                    "3. The app registration has 'Sites.Read.All' and 'Files.ReadWrite.All' permissions"
                )
        return await func(*args, **kwargs)
    return wrapper


def get_document_library_path() -> str:
    """Get the document library path from environment"""
    return os.getenv("SHP_DOC_LIBRARY", "Shared Documents")


@app.list_resources()
async def list_resources() -> list[Resource]:
    """List available SharePoint resources"""
    return [
        Resource(
            uri=AnyUrl(f"sharepoint:///{get_document_library_path()}"),
            name=f"SharePoint Document Library: {get_document_library_path()}",
            mimeType="application/vnd.sharepoint.folder",
            description="Main SharePoint document library configured for this server"
        )
    ]


@app.list_tools()
async def list_tools() -> list[Tool]:
    """List available SharePoint tools"""
    return [
        Tool(
            name="List_SharePoint_Folders",
            description="List all folders in a specified directory or root of the document library",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Path to the folder (relative to document library root). Leave empty for root.",
                        "default": ""
                    }
                }
            }
        ),
        Tool(
            name="List_SharePoint_Documents",
            description="List all documents in a specified folder with metadata",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Path to the folder containing documents",
                        "default": ""
                    }
                },
                "required": []
            }
        ),
        Tool(
            name="Get_Document_Content",
            description="Get the content of a document (supports text extraction from PDF, Word, Excel, and text files)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the file (relative to document library root)"
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="Upload_Document",
            description="Upload a new document to SharePoint",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Destination folder path"
                    },
                    "file_name": {
                        "type": "string",
                        "description": "Name of the file to create"
                    },
                    "content": {
                        "type": "string",
                        "description": "File content (text or base64 encoded for binary files)"
                    },
                    "is_binary": {
                        "type": "boolean",
                        "description": "Whether the content is base64 encoded binary",
                        "default": False
                    }
                },
                "required": ["folder_path", "file_name", "content"]
            }
        ),
        Tool(
            name="Update_Document",
            description="Update the content of an existing document",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the file to update"
                    },
                    "content": {
                        "type": "string",
                        "description": "New file content"
                    },
                    "is_binary": {
                        "type": "boolean",
                        "description": "Whether the content is base64 encoded binary",
                        "default": False
                    }
                },
                "required": ["file_path", "content"]
            }
        ),
        Tool(
            name="Delete_Document",
            description="Delete a document from SharePoint",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the file to delete"
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="Create_Folder",
            description="Create a new folder in SharePoint",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Path where to create the folder"
                    },
                    "folder_name": {
                        "type": "string",
                        "description": "Name of the new folder"
                    }
                },
                "required": ["folder_path", "folder_name"]
            }
        ),
        Tool(
            name="Delete_Folder",
            description="Delete an empty folder from SharePoint",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Path to the folder to delete"
                    }
                },
                "required": ["folder_path"]
            }
        ),
        Tool(
            name="Get_SharePoint_Tree",
            description="Get a recursive tree view of SharePoint folder structure",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Starting folder path (leave empty for root)",
                        "default": ""
                    },
                    "max_depth": {
                        "type": "integer",
                        "description": "Maximum depth to traverse",
                        "default": 5
                    }
                },
                "required": []
            }
        ),
        Tool(
            name="Test_Connection",
            description="Test the SharePoint connection and authentication",
            inputSchema={
                "type": "object",
                "properties": {},
                "required": []
            }
        ),
        Tool(
            name="Create_Word_Document",
            description="Create a formatted Word document (.docx) with headings, paragraphs, tables, lists, images, hyperlinks, horizontal rules, and page breaks. Supports rich formatting including colors, highlighting, and spacing. Upload directly to SharePoint.",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Destination folder path in SharePoint",
                        "default": ""
                    },
                    "file_name": {
                        "type": "string",
                        "description": "Name of the Word document (should end with .docx)"
                    },
                    "content": {
                        "type": "array",
                        "description": "Array of content items (headings, paragraphs, tables, lists, images, hyperlinks, horizontal rules, page breaks)",
                        "items": {
                            "type": "object",
                            "properties": {
                                "type": {
                                    "type": "string",
                                    "enum": ["heading", "paragraph", "table", "list", "image", "hyperlink", "horizontal_rule", "page_break"],
                                    "description": "Type of content element"
                                },
                                "level": {
                                    "type": "integer",
                                    "description": "Heading level (1-9, for type='heading')",
                                    "minimum": 1,
                                    "maximum": 9
                                },
                                "text": {
                                    "type": "string",
                                    "description": "Text content (for heading or paragraph)"
                                },
                                "bold": {
                                    "type": "boolean",
                                    "description": "Make text bold (for paragraph)"
                                },
                                "italic": {
                                    "type": "boolean",
                                    "description": "Make text italic (for paragraph)"
                                },
                                "underline": {
                                    "type": "boolean",
                                    "description": "Make text underlined (for paragraph)"
                                },
                                "font_size": {
                                    "type": "integer",
                                    "description": "Font size in points (for paragraph)"
                                },
                                "align": {
                                    "type": "string",
                                    "enum": ["left", "center", "right"],
                                    "description": "Text alignment (for paragraph)"
                                },
                                "color": {
                                    "description": "Text color - named color ('red', 'blue') or RGB array [255,0,0] (for paragraph)"
                                },
                                "highlight": {
                                    "type": "string",
                                    "enum": ["yellow", "bright_green", "turquoise", "pink", "blue", "red", "dark_blue", "dark_yellow", "green", "gray_25"],
                                    "description": "Highlight color (for paragraph)"
                                },
                                "line_spacing": {
                                    "type": "number",
                                    "description": "Line spacing multiplier, e.g., 1.5 for 1.5x spacing (for paragraph)"
                                },
                                "space_before": {
                                    "type": "number",
                                    "description": "Space before paragraph in points (for paragraph)"
                                },
                                "space_after": {
                                    "type": "number",
                                    "description": "Space after paragraph in points (for paragraph)"
                                },
                                "source": {
                                    "type": "string",
                                    "description": "Image source - URL (http/https) or base64 data (for type='image')"
                                },
                                "width": {
                                    "type": "number",
                                    "description": "Image width in inches (for type='image')"
                                },
                                "height": {
                                    "type": "number",
                                    "description": "Image height in inches (for type='image')"
                                },
                                "url": {
                                    "type": "string",
                                    "description": "Hyperlink URL (for type='hyperlink')"
                                },
                                "headers": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "Table column headers (for type='table')"
                                },
                                "rows": {
                                    "type": "array",
                                    "items": {
                                        "type": "array",
                                        "items": {"type": "string"}
                                    },
                                    "description": "Table rows data (for type='table')"
                                },
                                "style": {
                                    "type": "string",
                                    "description": "List style 'bullet' or 'number' (for type='list'), or table style name (for type='table')"
                                },
                                "items": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "List items (for type='list')"
                                }
                            }
                        }
                    }
                },
                "required": ["file_name", "content"]
            }
        ),
        Tool(
            name="Edit_Word_Document",
            description="""Edit an existing Word document (.docx) with find/replace operations and SECTION REPLACEMENT for monthly status reports. Preserves all formatting.

SIMPLE REPLACEMENT: Replace individual placeholders ({{company_name}} -> 'Acme Corp')

SECTION REPLACEMENT FOR MONTHLY REPORTS: Replace entire sections (8 bullets, tables) in one operation.
- Perfect for monthly status reports where sections need to be updated each month
- Use section markers: {{section:name}} ... {{/section:name}}
- Replace all content between markers with fresh bullet points or table rows
- Preserves formatting and list styles

USE CASES:
- Monthly status reports: Replace {{section:status_bullets}} with 8 new status items
- Metrics tables: Replace {{section:metrics}} with updated monthly data
- Template filling: Replace {{company_name}} and other placeholders with actual values""",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the existing Word document in SharePoint"
                    },
                    "replacements": {
                        "type": "array",
                        "description": """Array of replacement operations. Supports TWO FORMATS:

1. SIMPLE REPLACEMENT (single placeholder):
   {"find": "{{company_name}}", "replace": "Acme Corp"}

2. SECTION REPLACEMENT (for monthly reports - replaces entire sections like 8 bullets):
   {
     "find": "{{section:status_bullets}}",
     "replace_section": ["Item 1", "Item 2", "Item 3", ...],
     "section_type": "bullets"
   }

Template must have section markers:
{{section:status_bullets}}
* Old bullet 1
* Old bullet 2
{{/section:status_bullets}}

Section types: "bullets" (for bullet lists) or "table" (for tables)""",
                        "items": {
                            "type": "object",
                            "properties": {
                                "find": {
                                    "type": "string",
                                    "description": "Text to find. For simple: '{{company_name}}'. For sections: '{{section:status_bullets}}'"
                                },
                                "replace": {
                                    "type": "string",
                                    "description": "Text to replace with (for simple replacements only)"
                                },
                                "replace_section": {
                                    "type": "array",
                                    "description": "Array of items to replace section with. For bullets: array of strings. For tables: array of arrays.",
                                    "items": {}
                                },
                                "section_type": {
                                    "type": "string",
                                    "enum": ["bullets", "table"],
                                    "description": "Type of section (required for section replacements): 'bullets' for bullet lists, 'table' for tables"
                                }
                            },
                            "required": ["find"]
                        }
                    },
                    "output_file_path": {
                        "type": "string",
                        "description": "Optional: Path for the modified document. If not provided, overwrites the original file."
                    }
                },
                "required": ["file_path", "replacements"]
            }
        ),
        Tool(
            name="Create_PowerPoint",
            description="Create a formatted PowerPoint presentation (.pptx) with multiple slide layouts, text boxes, bullet points, images, tables, shapes, and rich formatting. Upload directly to SharePoint.",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_path": {
                        "type": "string",
                        "description": "Destination folder path in SharePoint",
                        "default": ""
                    },
                    "file_name": {
                        "type": "string",
                        "description": "Name of the PowerPoint file (should end with .pptx)"
                    },
                    "content": {
                        "type": "array",
                        "description": "Array of slide objects with layouts and content",
                        "items": {
                            "type": "object",
                            "properties": {
                                "type": {
                                    "type": "string",
                                    "enum": ["slide"],
                                    "description": "Must be 'slide'"
                                },
                                "layout": {
                                    "type": "string",
                                    "enum": ["title", "content", "section_header", "blank", "title_only"],
                                    "description": "Slide layout type",
                                    "default": "blank"
                                },
                                "title": {
                                    "type": "string",
                                    "description": "Slide title text"
                                },
                                "subtitle": {
                                    "type": "string",
                                    "description": "Subtitle text (for title/section_header layouts)"
                                },
                                "title_format": {
                                    "type": "object",
                                    "description": "Formatting for title (font_size, color, bold, italic, alignment)",
                                    "properties": {
                                        "font_size": {"type": "number"},
                                        "color": {"description": "Color as name or RGB array [255,0,0]"},
                                        "bold": {"type": "boolean"},
                                        "italic": {"type": "boolean"},
                                        "alignment": {"type": "string", "enum": ["left", "center", "right", "justify"]}
                                    }
                                },
                                "subtitle_format": {
                                    "type": "object",
                                    "description": "Formatting for subtitle"
                                },
                                "content": {
                                    "type": "array",
                                    "description": "Array of content items for the slide",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "type": {
                                                "type": "string",
                                                "enum": ["text_box", "bullets", "image", "table", "shape"],
                                                "description": "Content item type"
                                            },
                                            "text": {
                                                "type": "string",
                                                "description": "Text content (for text_box, shape)"
                                            },
                                            "left": {
                                                "type": "number",
                                                "description": "Left position in inches",
                                                "default": 1
                                            },
                                            "top": {
                                                "type": "number",
                                                "description": "Top position in inches",
                                                "default": 2
                                            },
                                            "width": {
                                                "type": "number",
                                                "description": "Width in inches",
                                                "default": 8
                                            },
                                            "height": {
                                                "type": "number",
                                                "description": "Height in inches"
                                            },
                                            "formatting": {
                                                "type": "object",
                                                "description": "Text formatting (font_size, color, bold, italic, alignment, vertical_alignment)",
                                                "properties": {
                                                    "font_size": {"type": "number"},
                                                    "color": {"description": "Color as name or RGB array"},
                                                    "bold": {"type": "boolean"},
                                                    "italic": {"type": "boolean"},
                                                    "alignment": {"type": "string", "enum": ["left", "center", "right", "justify"]},
                                                    "vertical_alignment": {"type": "string", "enum": ["top", "middle", "bottom"]}
                                                }
                                            },
                                            "items": {
                                                "type": "array",
                                                "items": {"type": "string"},
                                                "description": "Bullet point items (for bullets type)"
                                            },
                                            "ordered": {
                                                "type": "boolean",
                                                "description": "Use numbered list instead of bullets (for bullets type)",
                                                "default": False
                                            },
                                            "source": {
                                                "type": "string",
                                                "description": "Image source - URL or base64 (for image type)"
                                            },
                                            "rows": {
                                                "type": "array",
                                                "items": {
                                                    "type": "array",
                                                    "items": {"type": "string"}
                                                },
                                                "description": "Table rows including header row (for table type)"
                                            },
                                            "has_header": {
                                                "type": "boolean",
                                                "description": "First row is header (for table type)",
                                                "default": True
                                            },
                                            "header_color": {
                                                "description": "Header background color as name or RGB array (for table type)"
                                            },
                                            "shape_type": {
                                                "type": "string",
                                                "enum": ["rectangle", "rounded_rectangle", "oval", "diamond", "triangle", "arrow_right", "arrow_left", "star", "hexagon", "cloud"],
                                                "description": "Shape type (for shape type)",
                                                "default": "rectangle"
                                            },
                                            "fill_color": {
                                                "description": "Shape fill color as name or RGB array (for shape type)"
                                            },
                                            "line_color": {
                                                "description": "Shape line color (for shape type)"
                                            },
                                            "line_width": {
                                                "type": "number",
                                                "description": "Shape line width in points (for shape type)"
                                            },
                                            "text_format": {
                                                "type": "object",
                                                "description": "Text formatting for shape text (for shape type)"
                                            }
                                        }
                                    }
                                }
                            }
                        }
                    }
                },
                "required": ["file_name", "content"]
            }
        )
    ]


@app.call_tool()
@ensure_context
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    """Handle tool execution"""
    
    try:
        if name == "Test_Connection":
            return await test_connection()
        elif name == "List_SharePoint_Folders":
            return await list_folders(arguments.get("folder_path", ""))
        elif name == "List_SharePoint_Documents":
            return await list_documents(arguments.get("folder_path", ""))
        elif name == "Get_Document_Content":
            return await get_document_content(arguments["file_path"])
        elif name == "Upload_Document":
            return await upload_document(
                arguments["folder_path"],
                arguments["file_name"],
                arguments["content"],
                arguments.get("is_binary", False)
            )
        elif name == "Update_Document":
            return await update_document(
                arguments["file_path"],
                arguments["content"],
                arguments.get("is_binary", False)
            )
        elif name == "Delete_Document":
            return await delete_document(arguments["file_path"])
        elif name == "Create_Folder":
            return await create_folder(arguments["folder_path"], arguments["folder_name"])
        elif name == "Delete_Folder":
            return await delete_folder(arguments["folder_path"])
        elif name == "Get_SharePoint_Tree":
            return await get_tree(
                arguments.get("folder_path", ""),
                arguments.get("max_depth", 5)
            )
        elif name == "Create_Word_Document":
            return await create_word_document_tool(
                arguments.get("folder_path", ""),
                arguments["file_name"],
                arguments["content"]
            )
        elif name == "Edit_Word_Document":
            return await edit_word_document_tool(
                arguments["file_path"],
                arguments["replacements"],
                arguments.get("output_file_path")
            )
        elif name == "Create_PowerPoint":
            return await create_powerpoint_tool(
                arguments.get("folder_path", ""),
                arguments["file_name"],
                arguments["content"]
            )
        else:
            raise ValueError(f"Unknown tool: {name}")
            
    except Exception as e:
        logger.exception(f"Tool '{name}' failed")  # <-- prints stack trace
        return [TextContent(
            type="text",
            text=f"Error executing {name}: {str(e)}"
        )]



async def test_connection() -> list[TextContent]:
    """Test SharePoint connection using Microsoft Graph API"""
    try:
        logger.info("Testing SharePoint connection")
        # Test connection by fetching site ID and drive ID
        await asyncio.to_thread(graph_client._get_site_id)
        await asyncio.to_thread(graph_client._get_drive_id)
        logger.info("Connection test successful")

        return [TextContent(
            type="text",
            text=f"Successfully connected to SharePoint. Site URL: {graph_client.site_url}, Document Library: {graph_client.document_library_name}, Cloud: {graph_client.graph_endpoint}. Connection is working correctly."
        )]
    except Exception as e:
        logger.error(f"Connection test failed: {e}", exc_info=True)
        return [TextContent(
            type="text",
            text=f"Connection failed: {str(e)}. Common causes: 1. Incorrect credentials 2. Missing Microsoft Graph API permissions 3. Network/firewall issues. Set DEBUG=true for detailed diagnostics."
        )]


async def create_word_document_tool(folder_path: str, file_name: str, content: list) -> list[TextContent]:
    """Create a Word document and upload to SharePoint"""
    try:
        # Ensure filename ends with .docx
        if not file_name.lower().endswith('.docx'):
            file_name = f"{file_name}.docx"

        # Create Word document from JSON content
        logger.info(f"Creating Word document: {file_name}")
        doc_bytes = await asyncio.to_thread(create_word_document, content)
        logger.info(f"Document created successfully: {len(doc_bytes)} bytes")

        # Upload to SharePoint using Graph API
        logger.info(f"Uploading Word document to SharePoint: folder='{folder_path}', file='{file_name}'")
        try:
            await asyncio.to_thread(
                graph_client.upload_file,
                folder_path,
                file_name,
                doc_bytes
            )
            logger.info(f"Successfully uploaded Word document: {file_name}")
        except Exception as upload_error:
            logger.error(f"Failed to upload Word document: {upload_error}", exc_info=True)
            raise Exception(f"Document was created but upload failed: {str(upload_error)}")

        # Calculate statistics
        heading_count = len([i for i in content if i.get('type') == 'heading'])
        paragraph_count = len([i for i in content if i.get('type') == 'paragraph'])
        table_count = len([i for i in content if i.get('type') == 'table'])
        list_count = len([i for i in content if i.get('type') == 'list'])
        total_items = heading_count + paragraph_count + table_count + list_count

        # Simplified response to avoid OpenWebUI rendering issues
        return [TextContent(
            type="text",
            text=f"Successfully created and uploaded Word document '{file_name}' to '{folder_path or 'root'}'. "
                 f"Document contains {total_items} elements ({heading_count} headings, {paragraph_count} paragraphs, "
                 f"{table_count} tables, {list_count} lists)."
        )]

    except Exception as e:
        logger.error(f"Error creating Word document: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error creating Word document: {str(e)}")]


async def edit_word_document_tool(file_path: str, replacements: list, output_file_path: str = None) -> list[TextContent]:
    """Edit an existing Word document with find/replace operations"""
    try:
        # Download the existing document
        doc_bytes = await asyncio.to_thread(
            graph_client.get_file_content,
            file_path
        )

        # Perform find/replace operations
        modified_bytes, stats = await asyncio.to_thread(
            edit_word_document,
            doc_bytes,
            replacements
        )

        # Determine output path
        final_path = output_file_path if output_file_path else file_path

        # Extract folder and filename from path
        if '/' in final_path:
            parts = final_path.rsplit('/', 1)
            folder = parts[0]
            filename = parts[1]
        else:
            folder = ""
            filename = final_path

        # Ensure filename ends with .docx
        if not filename.lower().endswith('.docx'):
            filename = f"{filename}.docx"

        # Upload the modified document
        logger.info(f"Uploading edited document: folder='{folder}', filename='{filename}', size={len(modified_bytes)} bytes")
        try:
            await asyncio.to_thread(
                graph_client.upload_file,
                folder,
                filename,
                modified_bytes
            )
            logger.info(f"Successfully uploaded edited document to: {final_path}")
        except Exception as upload_error:
            logger.error(f"Failed to upload edited document: {upload_error}", exc_info=True)
            raise Exception(f"Document was edited successfully, but upload failed: {str(upload_error)}")

        # Build success message with statistics (simplified for OpenWebUI compatibility)
        total_replacements = stats.get('total_replacements', 0)
        paragraphs_modified = stats.get('paragraphs_modified', 0)
        tables_modified = stats.get('tables_modified', 0)
        sections_replaced = stats.get('sections_replaced', 0)

        return [TextContent(
            type="text",
            text=f"Successfully edited and uploaded Word document '{final_path}'. "
                 f"Made {total_replacements} replacements "
                 f"({paragraphs_modified} paragraphs, {tables_modified} tables, {sections_replaced} sections). "
                 f"All formatting preserved."
        )]

    except Exception as e:
        logger.error(f"Error editing Word document: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error editing Word document: {str(e)}")]


async def create_powerpoint_tool(folder_path: str, file_name: str, content: list) -> list[TextContent]:
    """Create a PowerPoint presentation and upload to SharePoint"""
    try:
        # Ensure filename ends with .pptx
        if not file_name.lower().endswith('.pptx'):
            file_name = f"{file_name}.pptx"

        # Create PowerPoint presentation from JSON content
        logger.info(f"Creating PowerPoint presentation: {file_name}")
        pptx_bytes = await asyncio.to_thread(create_powerpoint, content)
        logger.info(f"PowerPoint created successfully: {len(pptx_bytes)} bytes")

        # Upload to SharePoint using Graph API
        logger.info(f"Uploading PowerPoint to SharePoint: folder='{folder_path}', file='{file_name}'")
        try:
            await asyncio.to_thread(
                graph_client.upload_file,
                folder_path,
                file_name,
                pptx_bytes
            )
            logger.info(f"Successfully uploaded PowerPoint: {file_name}")
        except Exception as upload_error:
            logger.error(f"Failed to upload PowerPoint: {upload_error}", exc_info=True)
            raise Exception(f"Presentation was created but upload failed: {str(upload_error)}")

        # Calculate statistics
        slides_count = len([i for i in content if i.get('type') == 'slide'])

        # Count content items across all slides
        text_boxes = 0
        bullets = 0
        images = 0
        tables = 0
        shapes = 0

        for slide in content:
            if slide.get('type') == 'slide':
                for item in slide.get('content', []):
                    item_type = item.get('type', '')
                    if item_type == 'text_box':
                        text_boxes += 1
                    elif item_type == 'bullets':
                        bullets += len(item.get('items', []))
                    elif item_type == 'image':
                        images += 1
                    elif item_type == 'table':
                        tables += 1
                    elif item_type == 'shape':
                        shapes += 1

        return [TextContent(
            type="text",
            text=f"Successfully created and uploaded PowerPoint presentation '{file_name}' to '{folder_path or 'root'}'. "
                 f"Contains {slides_count} slides with {text_boxes} text boxes, {bullets} bullets, "
                 f"{images} images, {tables} tables, and {shapes} shapes."
        )]

    except Exception as e:
        logger.error(f"Error creating PowerPoint presentation: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error creating PowerPoint presentation: {str(e)}")]


async def list_folders(folder_path: str = "") -> list[TextContent]:
    """List folders in specified path using Microsoft Graph API"""
    doc_lib = get_document_library_path()
    full_path = f"{doc_lib}/{folder_path}" if folder_path else doc_lib

    try:
        logger.info(f"Listing folders in: {full_path}")
        # Use Graph API directly
        folders = await asyncio.to_thread(graph_client.list_folders, folder_path)
        folder_list = [f"{f['name']}" for f in folders]
        logger.info(f"Found {len(folders)} folders in: {full_path}")

        result = f"Folders in '{full_path}': " + ", ".join(folder_list) if folder_list else f"No folders found in '{full_path}'"
        return [TextContent(type="text", text=result)]

    except Exception as e:
        logger.error(f"Error listing folders: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error listing folders: {str(e)}")]


async def list_documents(folder_path: str = "") -> list[TextContent]:
    """List documents in specified folder using Microsoft Graph API"""
    doc_lib = get_document_library_path()
    full_path = f"{doc_lib}/{folder_path}" if folder_path else doc_lib

    try:
        logger.info(f"Listing documents in: {full_path}")
        # Use Graph API directly
        files = await asyncio.to_thread(graph_client.list_documents, folder_path)

        file_list = []
        for f in files:
            size_kb = f['size'] / 1024
            file_list.append(f"{f['name']} ({size_kb:.2f} KB)")
        logger.info(f"Found {len(files)} documents in: {full_path}")

        result = f"Documents in '{full_path}': " + ", ".join(file_list) if file_list else f"No documents found in '{full_path}'"
        return [TextContent(type="text", text=result)]

    except Exception as e:
        logger.error(f"Error listing documents: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error listing documents: {str(e)}")]


async def get_document_content(file_path: str) -> list[TextContent]:
    """Get document content using Microsoft Graph API"""
    try:
        logger.info(f"Getting document content: {file_path}")
        # Use Graph API to get file content
        content = await asyncio.to_thread(graph_client.get_file_content, file_path)
        logger.info(f"Retrieved document content: {file_path} ({len(content)} bytes)")

        ext = os.path.splitext(file_path)[1].lower()
        text_extensions = {'.txt', '.md', '.json', '.xml', '.html', '.csv', '.log'}
        excel_extensions = {'.xlsx', '.xlsm', '.xltx', '.xltm'}
        word_extensions = {'.docx'}
        powerpoint_extensions = {'.pptx'}

        # Handle text files
        if ext in text_extensions:
            text_content = content.decode("utf-8", errors="replace")
            return [TextContent(type="text", text=text_content)]

        # Handle Word documents
        if ext in word_extensions:
            try:
                from docx import Document

                # Load Word document from bytes
                word_file = BytesIO(content)
                doc = Document(word_file)

                result_text = f"Word document: {os.path.basename(file_path)}\n"
                result_text += f"Paragraphs: {len(doc.paragraphs)}\n"
                result_text += f"Tables: {len(doc.tables)}\n\n"
                result_text += "=== Content ===\n\n"

                # Extract text from paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():  # Only include non-empty paragraphs
                        result_text += para.text + "\n"

                # Extract text from tables
                if doc.tables:
                    result_text += "\n=== Tables ===\n\n"
                    for table_idx, table in enumerate(doc.tables, 1):
                        result_text += f"Table {table_idx}:\n"

                        # Get all rows
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            result_text += "  " + " | ".join(row_data) + "\n"

                        result_text += "\n"

                return [TextContent(type="text", text=result_text)]

            except Exception as word_error:
                logger.error(f"Error parsing Word document: {word_error}", exc_info=True)
                return [TextContent(
                    type="text",
                    text=f"Error parsing Word document: {str(word_error)}\n\n"
                         f"The file may be corrupted or in an unsupported format."
                )]

        # Handle PowerPoint presentations
        if ext in powerpoint_extensions:
            try:
                from pptx import Presentation

                # Load PowerPoint presentation from bytes
                pptx_file = BytesIO(content)
                prs = Presentation(pptx_file)

                result_text = f"PowerPoint presentation: {os.path.basename(file_path)}\n"
                result_text += f"Slides: {len(prs.slides)}\n\n"
                result_text += "=== Content ===\n\n"

                # Extract text from each slide
                for slide_idx, slide in enumerate(prs.slides, 1):
                    result_text += f"--- Slide {slide_idx} ---\n"

                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        # Handle shapes with text
                        if hasattr(shape, "text") and shape.text.strip():
                            result_text += shape.text.strip() + "\n"

                        # Handle tables
                        if hasattr(shape, "table"):
                            result_text += "\nTable:\n"
                            table = shape.table
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                result_text += "  " + " | ".join(row_data) + "\n"
                            result_text += "\n"

                    result_text += "\n"

                return [TextContent(type="text", text=result_text)]

            except Exception as pptx_error:
                logger.error(f"Error parsing PowerPoint presentation: {pptx_error}", exc_info=True)
                return [TextContent(
                    type="text",
                    text=f"Error parsing PowerPoint presentation: {str(pptx_error)}\n\n"
                         f"The file may be corrupted or in an unsupported format."
                )]

        # Handle Excel files
        if ext in excel_extensions:
            try:
                # Load Excel file from bytes
                excel_file = BytesIO(content)
                workbook = load_workbook(excel_file, data_only=True)

                result_text = f"Excel file: {os.path.basename(file_path)}\n"
                result_text += f"Sheets: {len(workbook.sheetnames)}\n\n"

                # Process each sheet
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    result_text += f"=== Sheet: {sheet_name} ===\n"
                    result_text += f"Dimensions: {sheet.dimensions}\n\n"

                    # Get all rows with data
                    rows_with_data = []
                    for row in sheet.iter_rows(values_only=True):
                        # Skip completely empty rows
                        if any(cell is not None for cell in row):
                            rows_with_data.append(row)

                    if rows_with_data:
                        # Calculate column widths for formatting
                        max_cols = max(len(row) for row in rows_with_data)
                        col_widths = [0] * max_cols

                        for row in rows_with_data:
                            for i, cell in enumerate(row):
                                if i < max_cols:
                                    cell_str = str(cell) if cell is not None else ""
                                    col_widths[i] = max(col_widths[i], len(cell_str))

                        # Display rows in tabular format
                        for row in rows_with_data:
                            row_data = []
                            for i, cell in enumerate(row):
                                if i < max_cols:
                                    cell_str = str(cell) if cell is not None else ""
                                    row_data.append(cell_str.ljust(col_widths[i]))
                            result_text += f"  {' | '.join(row_data)}\n"

                        result_text += f"\nTotal rows: {len(rows_with_data)}\n"
                    else:
                        result_text += "  (Sheet is empty)\n"

                    result_text += "\n"

                workbook.close()
                return [TextContent(type="text", text=result_text)]

            except Exception as excel_error:
                logger.error(f"Error parsing Excel file: {excel_error}", exc_info=True)
                return [TextContent(
                    type="text",
                    text=f"Error parsing Excel file: {str(excel_error)}\n\n"
                         f"The file may be corrupted or in an unsupported format."
                )]

        # Handle other binary files
        b64_content = base64.b64encode(content).decode("utf-8")
        return [TextContent(
            type="text",
            text=(
                "Binary file (base64 encoded):\n\n"
                f"{b64_content[:200]}...\n\n"
                f"Full content length: {len(b64_content)} characters"
            )
        )]

    except Exception as e:
        logger.error(f"Error reading document: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error reading document: {str(e)}")]


async def upload_document(folder_path: str, file_name: str, content: str, is_binary: bool = False) -> list[TextContent]:
    """Upload a document using Microsoft Graph API"""
    try:
        logger.info(f"Uploading document: folder='{folder_path}', file='{file_name}', is_binary={is_binary}")
        if is_binary:
            file_content = base64.b64decode(content)
        else:
            file_content = content.encode('utf-8')

        # Use Graph API to upload file
        try:
            await asyncio.to_thread(
                graph_client.upload_file,
                folder_path,
                file_name,
                file_content
            )
            logger.info(f"Successfully uploaded document: {file_name}")
        except Exception as upload_error:
            logger.error(f"Failed to upload document: {upload_error}", exc_info=True)
            raise Exception(f"Upload failed: {str(upload_error)}")

        return [TextContent(
            type="text",
            text=f"Successfully uploaded '{file_name}' to '{folder_path or 'root'}'"
        )]

    except Exception as e:
        logger.error(f"Error uploading document: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error uploading document: {str(e)}")]


async def update_document(file_path: str, content: str, is_binary: bool = False) -> list[TextContent]:
    """Update a document using Microsoft Graph API"""
    try:
        logger.info(f"Updating document: file='{file_path}', is_binary={is_binary}")
        if is_binary:
            file_content = base64.b64decode(content)
        else:
            file_content = content.encode('utf-8')

        # Split file_path into folder and filename
        folder_path = os.path.dirname(file_path)
        file_name = os.path.basename(file_path)

        # Use Graph API to upload/update file (PUT overwrites)
        try:
            await asyncio.to_thread(
                graph_client.upload_file,
                folder_path,
                file_name,
                file_content
            )
            logger.info(f"Successfully updated document: {file_path}")
        except Exception as upload_error:
            logger.error(f"Failed to update document: {upload_error}", exc_info=True)
            raise Exception(f"Update failed: {str(upload_error)}")

        return [TextContent(
            type="text",
            text=f"Successfully updated '{file_path}'"
        )]

    except Exception as e:
        logger.error(f"Error updating document: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error updating document: {str(e)}")]


async def delete_document(file_path: str) -> list[TextContent]:
    """Delete a document using Microsoft Graph API"""
    try:
        logger.info(f"Deleting document: {file_path}")
        # Use Graph API to delete file
        await asyncio.to_thread(graph_client.delete_file, file_path)
        logger.info(f"Successfully deleted document: {file_path}")

        return [TextContent(
            type="text",
            text=f"Successfully deleted '{file_path}'"
        )]

    except Exception as e:
        logger.error(f"Error deleting document: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error deleting document: {str(e)}")]


async def create_folder(folder_path: str, folder_name: str) -> list[TextContent]:
    """Create a folder using Microsoft Graph API"""
    try:
        logger.info(f"Creating folder: '{folder_name}' in '{folder_path or 'root'}'")
        # Use Graph API to create folder
        await asyncio.to_thread(
            graph_client.create_folder,
            folder_path,
            folder_name
        )
        logger.info(f"Successfully created folder: '{folder_name}' in '{folder_path or 'root'}'")

        return [TextContent(
            type="text",
            text=f"Successfully created folder '{folder_name}' in '{folder_path or 'root'}'"
        )]

    except Exception as e:
        logger.error(f"Error creating folder: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error creating folder: {str(e)}")]


async def delete_folder(folder_path: str) -> list[TextContent]:
    """Delete a folder using Microsoft Graph API"""
    try:
        logger.info(f"Deleting folder: {folder_path}")
        # Use Graph API to delete folder
        await asyncio.to_thread(graph_client.delete_folder, folder_path)
        logger.info(f"Successfully deleted folder: {folder_path}")

        return [TextContent(
            type="text",
            text=f"Successfully deleted folder '{folder_path}'"
        )]

    except Exception as e:
        logger.error(f"Error deleting folder: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error deleting folder: {str(e)}")]


async def get_tree(folder_path: str = "", max_depth: int = 5, current_depth: int = 0) -> list[TextContent]:
    """Get folder tree structure using Microsoft Graph API"""
    if current_depth >= max_depth:
        return [TextContent(type="text", text="Max depth reached")]

    try:
        if current_depth == 0:
            logger.info(f"Getting tree structure for: {folder_path or 'Root'}, max_depth={max_depth}")

        # Use Graph API to list folders
        folders = await asyncio.to_thread(graph_client.list_folders, folder_path)

        indent = "  " * current_depth
        tree_lines = [f"{indent}{folder_path or 'Root'}"]

        for f in folders:
            sub_path = f"{folder_path}/{f['name']}" if folder_path else f['name']
            sub_tree = await get_tree(sub_path, max_depth, current_depth + 1)
            tree_lines.append(sub_tree[0].text)

        if current_depth == 0:
            logger.info(f"Successfully generated tree structure for: {folder_path or 'Root'}")

        return [TextContent(type="text", text="\n".join(tree_lines))]

    except Exception as e:
        logger.error(f"Error getting tree: {e}", exc_info=True)
        return [TextContent(type="text", text=f"Error getting tree: {str(e)}")]


async def main():
    """Main entry point"""
    async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
        await app.run(
            read_stream,
            write_stream,
            app.create_initialization_options()
        )


if __name__ == "__main__":
    asyncio.run(main())

def run():
    """Sync entry point for the package"""
    asyncio.run(main())